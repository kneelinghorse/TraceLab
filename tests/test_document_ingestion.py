"""End-to-end ingestion pipeline tests covering prioritized formats."""

from __future__ import annotations

import json
from collections.abc import Callable
from pathlib import Path

import pytest

pytest.importorskip("docx")
pytest.importorskip("pptx")
pytest.importorskip("pandas")
pytest.importorskip("reportlab")

import pandas as pd
from docx import Document as DocxDocument
from httpx import ASGITransport, AsyncClient
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.pdfgen import canvas

import app.api.v1.documents as documents_api
from app.core.database import get_db
from app.main import app
from app.models.chunk import DocumentChunk
from app.models.document import Document
from app.models.processing_status import DocumentProcessingStatus
from app.services.chunking import ChunkingService
from app.services.coverage_report import CoverageReportGenerator
from app.services.document_ingestion import DocumentIngestionService
from app.services.processing_status import ProcessingStatusRecorder

LONG_TEXT = (
    "The quick brown fox jumps over the lazy dog near the UX research lab. " * 60
).strip()
COVERAGE_PATH = Path("cmos/reports/sprint-01/ingestion_format_coverage.json")


class StubRedactionService:
    """Lightweight stand-in for Presidio to keep tests deterministic."""

    def __init__(self, prefix: str = "[REDACTED]") -> None:
        self.prefix = prefix
        self.calls = []

    def redact_document(
        self, text: str, document_id: str, metadata=None, use_pseudonymization=True
    ):
        self.calls.append((text, document_id))
        redacted_text = f"{self.prefix} {text}"
        return {
            "redacted_text": redacted_text,
            "entities": [],
            "audit_trail": {
                "document_id": document_id,
                "metadata": metadata or {},
                "strategy": "stub",
            },
        }


def build_pdf(path: Path, text: str) -> Path:
    """Create a simple PDF using reportlab."""
    c = canvas.Canvas(str(path))
    text_object = c.beginText(40, 800)
    for line in text.split(". "):
        text_object.textLine(line.strip())
    c.drawText(text_object)
    c.save()
    return path


def build_docx(path: Path, text: str) -> Path:
    """Create DOCX with repeated paragraph."""
    doc = DocxDocument()
    doc.add_paragraph(text)
    doc.save(str(path))
    return path


def build_pptx(path: Path, text: str) -> Path:
    """Create PPTX slide with textbox content."""
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[5]  # blank layout
    slide = presentation.slides.add_slide(slide_layout)
    name_height = Inches(1.5)
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), name_height)
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.paragraphs[0].font.size = Pt(12)
    presentation.save(str(path))
    return path


def build_csv(path: Path, text: str) -> Path:
    """Create CSV with repeated rows."""
    sentences = [segment.strip() for segment in text.split(". ") if segment.strip()]
    df = pd.DataFrame(
        [
            {"question": f"Q{i + 1}", "response": sentence}
            for i, sentence in enumerate(sentences)
        ]
    )
    df.to_csv(path, index=False)
    return path


def build_xlsx(path: Path, text: str) -> Path:
    """Create XLSX workbook mirroring CSV structure."""
    sentences = [segment.strip() for segment in text.split(". ") if segment.strip()]
    df = pd.DataFrame(
        [
            {"question": f"Q{i + 1}", "response": sentence}
            for i, sentence in enumerate(sentences)
        ]
    )
    df.to_excel(path, index=False)
    return path


def build_markdown(path: Path, text: str) -> Path:
    """Create Markdown document with YAML front matter and rich headings."""
    headings = [
        f"## Insight {i + 1}\n\n{sentence.strip()}"
        for i, sentence in enumerate(text.split(". ")[:10], 1)
    ]
    body = "\n\n".join(headings) + "\n\n" + text
    payload = (
        "---\nproject_id: stub-project\ndoc_type: research_brief\n---\n\n# Markdown Ingestion Test\n\n"
        + body
    )
    path.write_text(payload)
    return path


FORMAT_BUILDERS: tuple[tuple[str, Callable[[Path, str], Path]], ...] = (
    ("pdf", build_pdf),
    ("docx", build_docx),
    ("pptx", build_pptx),
    ("csv", build_csv),
    ("xlsx", build_xlsx),
    ("md", build_markdown),
)


@pytest.mark.parametrize(
    "format_name,builder", FORMAT_BUILDERS, ids=[fmt for fmt, _ in FORMAT_BUILDERS]
)
def test_ingestion_pipeline_handles_all_formats(
    format_name: str, builder, db_session, project, tmp_path
):
    """Ensure each prioritized format runs through parsing, redaction, chunking, and audit logging."""
    file_path = builder(tmp_path / f"sample.{format_name}", LONG_TEXT)

    mime_types = {
        "pdf": "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "csv": "text/csv",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "md": "text/markdown",
    }
    file_types = {
        "pdf": "report",
        "docx": "report",
        "pptx": "report",
        "csv": "survey",
        "xlsx": "survey",
        "md": "notes",
    }

    document = Document(
        project_id=project.id,
        name=file_path.name,
        file_path=str(file_path),
        file_type=file_types.get(format_name, "report"),
        mime_type=mime_types[format_name],
        processed=False,
        chunked=False,
        validation_status="pending",
    )
    db_session.add(document)
    db_session.commit()
    db_session.refresh(document)

    # Record upload audit entry
    ProcessingStatusRecorder().record(
        db_session,
        document.id,
        stage="uploaded",
        status="succeeded",
        details={"format": format_name.upper(), "file_name": document.name},
    )

    ingestion_service = DocumentIngestionService(
        redaction_service=StubRedactionService(),
        chunking_service=ChunkingService(),
        status_recorder=ProcessingStatusRecorder(),
        coverage_report_generator=CoverageReportGenerator(),
    )

    result = ingestion_service.process_document(
        db=db_session, document_id=document.id, file_path=file_path
    )

    assert result["status"] == "completed"
    assert set(["extracted", "redacted", "chunked"]).issubset(result["stages"])

    # Reload document and related artifacts
    persisted = db_session.query(Document).filter(Document.id == document.id).one()
    assert persisted.processed is True
    assert persisted.chunked is True
    assert persisted.validation_status == "validated"
    assert persisted.content and persisted.content.startswith("[REDACTED]")

    chunks = (
        db_session.query(DocumentChunk)
        .filter(DocumentChunk.document_id == document.id)
        .all()
    )
    assert len(chunks) >= 1
    assert all(chunk.token_count >= 500 for chunk in chunks)

    statuses = (
        db_session.query(DocumentProcessingStatus)
        .filter(DocumentProcessingStatus.document_id == document.id)
        .all()
    )
    succeeded_stages = {(status.stage, status.status) for status in statuses}
    for stage in ("uploaded", "extracted", "redacted", "chunked"):
        assert (stage, "succeeded") in succeeded_stages

    # Coverage artifact should include the processed format
    assert COVERAGE_PATH.exists()
    payload = json.loads(COVERAGE_PATH.read_text())
    format_key = (
        persisted.file_path and Path(persisted.file_path).suffix.lstrip(".").upper()
    )
    assert format_key in payload["formats"]
    assert payload["formats"][format_key]["processed"] >= 1
    assert payload["formats"][format_key]["chunked"] >= 1


@pytest.mark.asyncio
async def test_upload_and_process_endpoints_record_statuses(
    db_session, project, tmp_path, monkeypatch, auth_headers
):
    """Verify FastAPI endpoints create audit events and trigger the ingestion pipeline."""

    stub_service = DocumentIngestionService(
        redaction_service=StubRedactionService(),
        chunking_service=ChunkingService(),
        status_recorder=ProcessingStatusRecorder(),
        coverage_report_generator=CoverageReportGenerator(),
    )

    monkeypatch.setattr(documents_api, "_ingestion_service", stub_service)
    monkeypatch.setattr(documents_api, "_ingestion_init_error", None)
    monkeypatch.setattr(documents_api, "get_ingestion_service", lambda: stub_service)

    file_path = build_csv(tmp_path / "api_upload.csv", LONG_TEXT)

    def override_get_db():
        yield db_session

    app.dependency_overrides[get_db] = override_get_db

    try:
        transport = ASGITransport(app=app)
        async with AsyncClient(
            transport=transport, base_url="http://test", headers=auth_headers
        ) as client:
            upload_response = await client.post(
                f"{documents_api.settings.api_v1_prefix}/documents/upload",
                params={"project_id": str(project.id)},
                files={"file": (file_path.name, file_path.read_bytes(), "text/csv")},
            )

            assert upload_response.status_code == 200
            uploaded_payload = upload_response.json()
            assert any(
                event["stage"] == "uploaded"
                for event in uploaded_payload["processing_events"]
            )

            document_id = uploaded_payload["id"]

            process_response = await client.post(
                f"{documents_api.settings.api_v1_prefix}/documents/{document_id}/process"
            )
            assert process_response.status_code == 200
            process_payload = process_response.json()
            assert process_payload["status"] == "completed"

            detail_response = await client.get(
                f"{documents_api.settings.api_v1_prefix}/documents/{document_id}"
            )
            assert detail_response.status_code == 200
            detail_payload = detail_response.json()
            succeeded = {
                event["stage"]
                for event in detail_payload["processing_events"]
                if event["status"] == "succeeded"
            }
            assert {"uploaded", "extracted", "redacted", "chunked"}.issubset(succeeded)

        assert COVERAGE_PATH.exists()
    finally:
        app.dependency_overrides.pop(get_db, None)
