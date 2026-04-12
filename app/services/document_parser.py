"""
Document parsers for various file formats.

Supports PDF, DOCX, PPTX, CSV, and XLSX formats as specified in Sprint 1 scope.
"""

import io
from pathlib import Path

try:
    from pdfminer.high_level import extract_text as extract_pdf_text

    PDFMINER_AVAILABLE = True
except ImportError:
    try:
        # Fallback: try pdfminer.six module structure
        from pdfminer.six.high_level import extract_text as extract_pdf_text

        PDFMINER_AVAILABLE = True
    except ImportError:
        PDFMINER_AVAILABLE = False
        extract_pdf_text = None

try:
    from docx import Document as DocxDocument

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import pandas as pd

    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False

import json
import xml.etree.ElementTree as ET

import yaml


class DocumentParser:
    """Parser for extracting text from various document formats."""

    @staticmethod
    def parse(file_path: Path, file_content: bytes | None = None) -> str:
        """
        Parse a document file and extract text content.

        Args:
            file_path: Path to the file (for format detection)
            file_content: Optional file content bytes (if file is already in memory)

        Returns:
            Extracted text content as string

        Raises:
            ValueError: If file format is not supported
            Exception: If parsing fails
        """
        suffix = file_path.suffix.lower()

        # Use file_content if provided, otherwise read from file_path
        if file_content is None:
            file_content = file_path.read_bytes()

        if suffix == ".pdf":
            return DocumentParser._parse_pdf(file_content)
        elif suffix == ".docx":
            return DocumentParser._parse_docx(file_content)
        elif suffix == ".pptx":
            return DocumentParser._parse_pptx(file_content)
        elif suffix == ".csv":
            return DocumentParser._parse_csv(file_content)
        elif suffix == ".xlsx":
            return DocumentParser._parse_xlsx(file_content)
        elif suffix in {".md", ".markdown", ".txt"}:
            return DocumentParser._parse_markdown(file_content)
        elif suffix == ".json":
            return DocumentParser._parse_json(file_content)
        elif suffix == ".xml":
            return DocumentParser._parse_xml(file_content)
        elif suffix in {".yaml", ".yml"}:
            return DocumentParser._parse_yaml(file_content)
        else:
            raise ValueError(f"Unsupported file format: {suffix}")

    @staticmethod
    def _parse_pdf(content: bytes) -> str:
        """Extract text from PDF using pdfminer.six."""
        if not PDFMINER_AVAILABLE:
            raise ImportError("pdfminer.six is not installed")

        # pdfminer expects a file path, so we write to a temporary buffer
        with io.BytesIO(content) as buffer:
            # Create a temporary file path for pdfminer
            import tempfile

            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp_file:
                tmp_file.write(content)
                tmp_path = tmp_file.name

            try:
                text = extract_pdf_text(tmp_path)
                return text if text else ""
            finally:
                Path(tmp_path).unlink(missing_ok=True)

    @staticmethod
    def _parse_docx(content: bytes) -> str:
        """Extract text from DOCX using python-docx."""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx is not installed")

        doc = DocxDocument(io.BytesIO(content))
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    paragraphs.append(" | ".join(row_text))

        return "\n".join(paragraphs)

    @staticmethod
    def _parse_pptx(content: bytes) -> str:
        """Extract text from PPTX using python-pptx."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is not installed")

        prs = Presentation(io.BytesIO(content))
        texts = []

        for slide in prs.slides:
            slide_texts = []

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

                # Extract text from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_texts.append(" | ".join(row_text))

            if slide_texts:
                texts.append("\n".join(slide_texts))

            # Extract speaker notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        texts.append(f"[Notes]\n{notes_text}")

        return "\n\n".join(texts)

    @staticmethod
    def _parse_csv(content: bytes) -> str:
        """Extract text from CSV using pandas."""
        if not PANDAS_AVAILABLE:
            raise ImportError("pandas is not installed")

        try:
            # Try UTF-8 first
            df = pd.read_csv(io.BytesIO(content), encoding="utf-8")
        except UnicodeDecodeError:
            # Fallback to latin-1 for broader compatibility
            df = pd.read_csv(io.BytesIO(content), encoding="latin-1")

        # Convert DataFrame to text representation
        return df.to_string(index=False)

    @staticmethod
    def _parse_xlsx(content: bytes) -> str:
        """Extract text from XLSX using pandas with openpyxl engine."""
        if not PANDAS_AVAILABLE:
            raise ImportError("pandas is not installed")

        # Read all sheets
        excel_file = pd.ExcelFile(io.BytesIO(content), engine="openpyxl")
        sheets_text = []

        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name, engine="openpyxl")
            sheet_content = f"[Sheet: {sheet_name}]\n{df.to_string(index=False)}"
            sheets_text.append(sheet_content)

        return "\n\n".join(sheets_text)

    @staticmethod
    def is_format_supported(file_path: Path) -> bool:
        """Check if file format is supported."""
        if file_path is None:
            return False
        suffix = file_path.suffix.lower()
        return suffix in {
            ".pdf",
            ".docx",
            ".pptx",
            ".csv",
            ".xlsx",
            ".md",
            ".markdown",
            ".txt",
            ".json",
            ".xml",
            ".yaml",
            ".yml",
        }

    @staticmethod
    def _parse_markdown(content: bytes) -> str:
        """Return Markdown content as plain text while stripping YAML front matter."""
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("utf-8", errors="ignore")

        stripped = text.lstrip()
        if stripped.startswith("---"):
            # YAML front matter ends at the next line starting with ---.
            front_matter_end = stripped.find("\n---", 3)
            if front_matter_end != -1:
                candidate = stripped[front_matter_end + 4 :]
                # Handle optional YAML document terminator ("...").
                candidate_stripped = candidate.lstrip()
                if candidate_stripped.startswith("..."):
                    terminator_end = candidate_stripped.find("\n")
                    candidate = (
                        candidate_stripped[terminator_end + 1 :]
                        if terminator_end != -1
                        else ""
                    )
                text = candidate

        return text.strip()

    @staticmethod
    def _parse_json(content: bytes) -> str:
        """Extract text from JSON file.

        Parses the JSON and pretty-prints it for readability.
        For arrays of objects, extracts text values recursively.
        """
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("utf-8", errors="ignore")

        try:
            data = json.loads(text)
            # Pretty print for structured viewing
            return json.dumps(data, indent=2, ensure_ascii=False)
        except json.JSONDecodeError as e:
            # Return raw content if parsing fails
            return f"[JSON Parse Error: {e}]\n\n{text}"

    @staticmethod
    def _parse_xml(content: bytes) -> str:
        """Extract text from XML file.

        Parses XML and extracts all text content, preserving structure
        through indentation.
        """
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("utf-8", errors="ignore")

        def extract_text(element: ET.Element, depth: int = 0) -> list[str]:
            """Recursively extract text from XML element."""
            result = []
            indent = "  " * depth
            tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

            # Handle element text
            if element.text and element.text.strip():
                result.append(f"{indent}<{tag}> {element.text.strip()}")
            elif len(element) == 0:  # Leaf element
                result.append(f"{indent}<{tag}>")
            else:
                result.append(f"{indent}<{tag}>")

            # Process children
            for child in element:
                result.extend(extract_text(child, depth + 1))

            # Handle tail text
            if element.tail and element.tail.strip():
                result.append(f"{indent}{element.tail.strip()}")

            return result

        try:
            root = ET.fromstring(text)
            lines = extract_text(root)
            return "\n".join(lines)
        except ET.ParseError as e:
            # Return raw content if parsing fails
            return f"[XML Parse Error: {e}]\n\n{text}"

    @staticmethod
    def _parse_yaml(content: bytes) -> str:
        """Extract text from YAML file.

        Parses YAML and pretty-prints it. Handles multi-document YAML files.
        """
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("utf-8", errors="ignore")

        try:
            # Handle multi-document YAML
            documents = list(yaml.safe_load_all(text))
            if len(documents) == 1:
                return yaml.dump(
                    documents[0], default_flow_style=False, allow_unicode=True
                )
            else:
                parts = []
                for i, doc in enumerate(documents):
                    if doc is not None:
                        parts.append(f"--- Document {i + 1} ---")
                        parts.append(
                            yaml.dump(doc, default_flow_style=False, allow_unicode=True)
                        )
                return "\n".join(parts)
        except yaml.YAMLError as e:
            # Return raw content if parsing fails
            return f"[YAML Parse Error: {e}]\n\n{text}"
